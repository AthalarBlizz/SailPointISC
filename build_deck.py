#!/usr/bin/env python3
"""Build DevDays 2026 presentation from SailPoint template."""

import zipfile, os, shutil
from pptx import Presentation
from pptx.util import Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from lxml import etree

TEMPLATE = os.path.expanduser(
    "~/work/ai/professional/devdays/INBOX/SailPoint-DeveloperDays-Template-2026.pptx"
)
OUTPUT = os.path.expanduser(
    "~/work/ai/professional/devdays/devdays2026.pptx"
)

# Layout indices
LAYOUT_TITLE_SINGLE  = 0   # Title - Single Speaker - Developer Days
LAYOUT_TITLE_CONTENT = 4   # Title and Content
LAYOUT_TWO_CONTENT   = 7   # Two Content
LAYOUT_SECTION       = 11  # Section Header DevDays
LAYOUT_THREE_COLUMNS = 6   # Three Columns
LAYOUT_BLANK         = 13  # Blank


def add_slide(prs, layout_idx):
    layout = prs.slide_layouts[layout_idx]
    return prs.slides.add_slide(layout)


def set_placeholder(slide, idx, text, font_size=None, bold=None):
    for shape in slide.placeholders:
        if shape.placeholder_format.idx == idx:
            tf = shape.text_frame
            tf.word_wrap = True
            tf.clear()
            p = tf.paragraphs[0]
            run = p.add_run()
            run.text = text
            if font_size:
                run.font.size = Pt(font_size)
            if bold is not None:
                run.font.bold = bold
            return shape
    return None


def set_body_bullets(slide, idx, bullets, font_size=18):
    """Set body placeholder with bullet list. bullets = list of (text, level, bold)."""
    for shape in slide.placeholders:
        if shape.placeholder_format.idx == idx:
            tf = shape.text_frame
            tf.word_wrap = True
            tf.clear()
            for i, item in enumerate(bullets):
                if isinstance(item, str):
                    text, level, bold = item, 0, False
                else:
                    text = item[0]
                    level = item[1] if len(item) > 1 else 0
                    bold = item[2] if len(item) > 2 else False

                if i == 0:
                    p = tf.paragraphs[0]
                else:
                    p = tf.add_paragraph()

                p.level = level
                run = p.add_run()
                run.text = text
                run.font.size = Pt(font_size)
                run.font.bold = bold
            return shape
    return None


def delete_template_slides(prs):
    """Remove all existing example slides from the template."""
    from pptx.oxml.ns import qn
    sldIdLst = prs.slides._sldIdLst
    for sldId in list(sldIdLst):
        rId = sldId.get(qn('r:id'))
        if rId:
            try:
                prs.part.drop_rel(rId)
            except Exception:
                pass
        sldIdLst.remove(sldId)


def save_clean(prs, path):
    """Save and deduplicate ZIP entries to prevent corruption."""
    tmp_raw = path + ".tmp_raw"
    tmp_clean = path + ".tmp_clean"
    prs.save(tmp_raw)
    seen = set()
    with zipfile.ZipFile(tmp_raw, 'r') as zin:
        with zipfile.ZipFile(tmp_clean, 'w', zipfile.ZIP_DEFLATED) as zout:
            for item in zin.infolist():
                if item.filename not in seen:
                    seen.add(item.filename)
                    zout.writestr(item, zin.read(item.filename))
    shutil.move(tmp_clean, path)
    os.remove(tmp_raw)


# ── Build ──────────────────────────────────────────────────────────────────

prs = Presentation(TEMPLATE)
delete_template_slides(prs)

# ── Slide 1: Title ──────────────────────────────────────────────────────────
s = add_slide(prs, LAYOUT_TITLE_SINGLE)
set_placeholder(s, 0, "When There's No Button for That\nAI-Assisted Development", font_size=36)
set_placeholder(s, 11,
    "speaker: Alicia Gutschow;\n"
    "title: Principal Solution Advisor;\n"
    "company: MajorKey Technologies;"
)

# ── Slide 2: There's Always a Gap ──────────────────────────────────────────
s = add_slide(prs, LAYOUT_TITLE_CONTENT)
set_placeholder(s, 0, "There's Always a Gap")
set_body_bullets(s, 1, [
    ("ISC handles a lot — but every program eventually hits the edge of what the platform does natively.", 0, False),
    ("", 0, False),
    ("Real scenarios that require API solutions:", 0, True),
    ("An external system needs to disable a user when something happens that ISC doesn't track", 1, False),
    ("Three systems each have part of an identity record, and a fourth system needs all of it", 1, False),
    ('A manager says "just give them what Jenny has" — but ISC access requests don\'t work that way', 1, False),
    ("", 0, False),
    ("The API surface exists to solve these problems. The challenge has always been the cost — navigating hundreds of endpoints, writing and testing code, and doing it again for the next use case.", 0, False),
    ("", 0, False),
    ("AI changes that equation.", 0, True),
], font_size=16)

# ── Slide 3: AI as a Force Multiplier ──────────────────────────────────────
s = add_slide(prs, LAYOUT_TITLE_CONTENT)
set_placeholder(s, 0, "AI as a Force Multiplier")
set_body_bullets(s, 1, [
    ("Without AI, solving an ISC API problem requires:", 0, True),
    ("Knowing what you are trying to accomplish — technically, not just conceptually", 1, False),
    ("Finding the right APIs, understanding the parameters, and knowing what to expect back", 1, False),
    ("Setting up a development environment and writing working code", 1, False),
    ("Executing code, interpreting results, and debugging when something goes wrong", 1, False),
    ("", 0, False),
    ("AI lets you start with plain English. It walks you through every one of those steps — teaching you what you want to learn and automating what you don't.", 0, False),
    ("", 0, False),
    ("It also works the other way: you inherited code from the previous engineer or a partner delivers code, and you want to know what the code actually does.", 0, False),
], font_size=16)

# ── Slide 4: Setting Up Your Environment ───────────────────────────────────
s = add_slide(prs, LAYOUT_THREE_COLUMNS)
set_placeholder(s, 0, "Setting Up Your Environment")
set_body_bullets(s, 1, [
    ("Python", 0, True),
    ("python.org — download for your platform", 1, False),
    ("Version 3.10 or higher", 1, False),
    ("Verify: python --version", 1, False),
], font_size=15)
set_body_bullets(s, 14, [("SailPoint Python SDK", 0, True)], font_size=15)
set_body_bullets(s, 15, [
    ("Wraps the ISC API — handles auth, pagination, and token refresh", 1, False),
    ("Install: pip install sailpoint", 1, False),
    ("Source: github.com/sailpoint-oss/python-sdk", 1, False),
], font_size=15)
set_body_bullets(s, 16, [("AI Coding Tool", 0, True)], font_size=15)
set_body_bullets(s, 17, [
    ("Claude Code, VS Code + Copilot, Cursor, or any AI-assisted environment", 1, False),
    ("This session uses Claude Code — the approach works with any of them", 1, False),
], font_size=15)

# ── Slide 5: Connecting to Your Tenant ─────────────────────────────────────
s = add_slide(prs, LAYOUT_TITLE_CONTENT)
set_placeholder(s, 0, "Connecting to Your Tenant")
set_body_bullets(s, 1, [
    ("Personal Access Token (PAT)", 0, True),
    ("ISC: top-right menu > Preferences > Personal Access Tokens", 1, False),
    ("Give it a descriptive name — you will want to know what it is for later", 1, False),
    ("Copy it immediately — ISC will not show it again", 1, False),
    ("", 0, False),
    ("Storing credentials with python-keyring", 0, True),
    ("Stores in your OS keystore — nothing written to disk", 1, False),
    ("macOS Keychain, Windows Credential Manager, Linux Secret Service", 2, False),
    ("Store once interactively, retrieve in every script", 1, False),
    ("pip install keyring", 1, False),
    ("", 0, False),
    ("For production environments", 0, True),
    ("Use what your team already has: AWS Secrets Manager, Azure Key Vault, CyberArk, or other secrets management", 1, False),
    ("", 0, False),
    ("What not to do", 0, True),
    ("Never put a token in plaintext in your code or a file you might share", 1, False),
    ("Never paste your token into an AI tool — the AI does not need your credentials to help you write code", 1, False),
], font_size=14)

# ── Slide 6: The SailPoint API and SDK ─────────────────────────────────────
s = add_slide(prs, LAYOUT_TITLE_CONTENT)
set_placeholder(s, 0, "The SailPoint API and SDK")
set_body_bullets(s, 1, [
    ("The OSS API specs — start here", 0, True),
    ("Complete OpenAPI specification for every ISC endpoint: github.com/sailpoint-oss/api-specs", 1, False),
    ("614 endpoint YAML files, 500+ data model schemas", 1, False),
    ("When AI reads these directly, it knows every parameter, required scope, and response field", 1, False),
    ("No documentation hunting — no hallucinated endpoints", 1, False),
    ("", 0, False),
    ("The API surface", 0, True),
    ("600+ REST API endpoints across versions", 1, False),
    ("Covers everything in the UI and significantly more", 1, False),
    ("Documentation: developer.sailpoint.com/docs/api/", 1, False),
    ("", 0, False),
    ("The SDK", 0, True),
    ("SailPoint-maintained libraries — available for Golang, PowerShell, and Python", 1, False),
    ("Handles OAuth2 token management, pagination, and retry automatically", 1, False),
], font_size=14)

# ── Slide 7: Before You Start ──────────────────────────────────────────────
s = add_slide(prs, LAYOUT_TITLE_CONTENT)
set_placeholder(s, 0, "Before You Start — What to Know")
set_body_bullets(s, 1, [
    ("Understand your AI tool's data handling", 0, True),
    ("When you interact with any AI tool, you are potentially exposing data", 1, False),
    ("Review your AI tool's data retention and training policies", 1, False),
    ("Consider the sensitivity of data required for your project", 1, False),
    ("Enterprise AI subscriptions with restricted data clauses", 2, False),
    ("AWS Bedrock or Azure OpenAI — same AI capability, data stays within your org's environment", 2, False),
    ("Always use your organization's approved AI platform if designated", 2, False),
    ("", 0, False),
    ("Least privilege for your credentials", 0, True),
    ("Scope your PAT and OAuth client to only the API rights the integration needs", 1, False),
    ("Never put credentials in code — always use environment variables or a secrets vault", 1, False),
    ("Ensure that protected data elements, such as IP addresses and working data exports, are never committed to source code repositories", 1, False),
], font_size=14)

# ── Slide 6: Closing (optional) ────────────────────────────────────────────
s = add_slide(prs, LAYOUT_SECTION)
set_placeholder(s, 0, "Thank you!")
set_placeholder(s, 1,
    "Code and starter kit: github.com/agutschow/2026DeveloperDays\n"
    "alicia.gutschow@majorkey.com"
)

# ── Save ───────────────────────────────────────────────────────────────────
save_clean(prs, OUTPUT)
print(f"Saved: {OUTPUT}")
